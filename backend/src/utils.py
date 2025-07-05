import os
import pandas as pd
from io import BytesIO
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pypdf
import docx
import pptx


def parse_pdf(file: BytesIO) -> str:
    pdf_reader = pypdf.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def parse_docx(file: BytesIO) -> str:
    document = docx.Document(file)
    return "\n".join([para.text for para in document.paragraphs])

def parse_pptx(file: BytesIO) -> str:
    presentation = pptx.Presentation(file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def parse_csv(file: BytesIO) -> str:
    df = pd.read_csv(file)
    return df.to_string()

def parse_txt(file: BytesIO) -> str:
    return file.read().decode("utf-8")

def parse_document(uploaded_file) -> str:
    """Parses an uploaded file based on its filename and content."""
    file_content = BytesIO(uploaded_file.file)
    file_extension = os.path.splitext(uploaded_file.filename)[1].lower()

    if file_extension == ".pdf":
        return parse_pdf(file_content)
    elif file_extension == ".docx":
        return parse_docx(file_content)
    elif file_extension == ".pptx":
        return parse_pptx(file_content)
    elif file_extension == ".csv":
        return parse_csv(file_content)
    elif file_extension in [".txt", ".md"]:
        return parse_txt(file_content)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

def get_text_chunks(text: str) -> list[str]:
    """Splits a long text into smaller, semantically meaningful chunks."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    return text_splitter.split_text(text)